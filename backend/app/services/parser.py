import fitz  # pymupdf
import pandas as pd
from pptx import Presentation
import pytesseract
from typing import Any
import os
import tempfile
import requests
from chromadb import PersistentClient

# Use the same ChromaDB setup as the generator
CHROMA_PATH = "./chromadb_reports"
vector_db = PersistentClient(path=CHROMA_PATH)

# Deepgram API configuration
DEEPGRAM_API_KEY = "e9660a02522e2ecf5f269805fdb3bab3dd31488a"

def extract_pdf(file_path):
    doc = fitz.open(file_path)
    texts = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text()
        texts.append({"page": page_num, "text": text})
    return texts

def extract_pptx(file_path):
    prs = Presentation(file_path)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slides.append({"slide": idx, "text": text})
    return slides

def extract_xlsx(file_path):
    df = pd.read_excel(file_path)
    rows = df.to_dict(orient="records")
    return [{"row": idx+1, "data": row} for idx, row in enumerate(rows)]

def extract_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return [{"text": f.read()}]

def extract_image(file_path):
    text = pytesseract.image_to_string(file_path)
    return [{"text": text}]

def extract_docx(file_path):
    """
    Extract text from Word documents (.doc, .docx)
    """
    try:
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract DOCX: {str(e)}"}]

def extract_archive(file_path):
    """
    Extract text from archive files (zip, rar, etc.)
    """
    try:
        import zipfile
        import tarfile
        import rarfile
        
        text_content = []
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == ".zip":
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                for file_info in zip_ref.filelist:
                    if not file_info.is_dir() and file_info.filename.endswith(('.txt', '.pdf', '.docx')):
                        try:
                            content = zip_ref.read(file_info.filename)
                            if file_info.filename.endswith('.txt'):
                                text_content.append(f"[From {file_info.filename}]: {content.decode('utf-8', errors='ignore')}")
                        except:
                            continue
        elif file_ext in [".tar", ".gz"]:
            with tarfile.open(file_path, 'r:*') as tar_ref:
                for member in tar_ref.getmembers():
                    if member.isfile() and member.name.endswith(('.txt', '.pdf', '.docx')):
                        try:
                            content = tar_ref.extractfile(member).read()
                            if member.name.endswith('.txt'):
                                text_content.append(f"[From {member.name}]: {content.decode('utf-8', errors='ignore')}")
                        except:
                            continue
        
        if text_content:
            return [{"text": "\n\n".join(text_content)}]
        else:
            return [{"text": f"Archive {os.path.basename(file_path)} contains no extractable text files"}]
            
    except Exception as e:
        return [{"error": f"Failed to extract archive: {str(e)}"}]

def extract_video_audio(file_path):
    """
    Extract transcript from video/audio files using Deepgram API
    """
    try:
        # Determine content type based on file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext in ['.wav', '.mp3', '.m4a', '.flac']:
            content_type = "audio/wav" if file_ext == '.wav' else "audio/mp4"
        elif file_ext in ['.mp4', '.avi', '.mov', '.mkv']:
            content_type = "video/mp4"
        else:
            content_type = "video/mp4"  # Default for video files
        
        # Open the file and send to Deepgram
        with open(file_path, "rb") as f:
            response = requests.post(
                "https://api.deepgram.com/v1/listen?punctuate=true&diarize=true&model=nova-2",
                headers={
                    "Authorization": f"Token {DEEPGRAM_API_KEY}",
                    "Content-Type": content_type,
                },
                data=f,
                timeout=300  # 5 minute timeout for large files
            )
        
        if response.status_code != 200:
            raise Exception(f"Deepgram API error: {response.status_code} - {response.text}")
        
        result = response.json()
        
        # Extract transcript
        transcript = result["results"]["channels"][0]["alternatives"][0]["transcript"]
        
        # Extract word-level timestamps for better context
        words = result["results"]["channels"][0]["alternatives"][0]["words"]
        
        # Create structured data with timestamps
        # Split transcript into meaningful chunks (sentences or phrases)
        import re
        sentences = re.split(r'[.!?]+', transcript)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        # Create chunks of 2-3 sentences each to avoid too many citations
        chunk_size = 2
        chunks = []
        for i in range(0, len(sentences), chunk_size):
            chunk_sentences = sentences[i:i + chunk_size]
            chunk_text = '. '.join(chunk_sentences)
            if chunk_text and not chunk_text.endswith('.'):
                chunk_text += '.'
            
            chunks.append({
                "text": chunk_text,
                "word_count": len(chunk_text.split()),
                "duration": words[-1]["end"] if words else 0,
                "words": words,
                "chunk_index": len(chunks) + 1
            })
        
        return chunks if chunks else [{
            "text": transcript,
            "word_count": len(words),
            "duration": words[-1]["end"] if words else 0,
            "words": words
        }]
        
    except Exception as e:
        print(f"Error transcribing video/audio: {str(e)}")
        return [{"error": f"Transcription failed: {str(e)}"}]

async def parse_file(file) -> dict:
    ext = os.path.splitext(file.filename)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name
    
    # Handle different file types
    if ext == ".pdf":
        data = extract_pdf(tmp_path)
    elif ext in [".ppt", ".pptx"]:
        data = extract_pptx(tmp_path)
    elif ext in [".xls", ".xlsx"]:
        data = extract_xlsx(tmp_path)
    elif ext in [".txt"]:
        data = extract_txt(tmp_path)
    elif ext in [".png", ".jpg", ".jpeg"]:
        data = extract_image(tmp_path)
    elif ext in [".mp4", ".avi", ".mov", ".mkv", ".wav", ".mp3", ".m4a", ".flac"]:
        data = extract_video_audio(tmp_path)
    else:
        data = [{"error": f"Unsupported file type: {ext}"}]
    os.remove(tmp_path)
    # Store in vector DB using the same collection as generator
    import uuid
    collection = vector_db.get_or_create_collection("reports")
    
    # Process and store the extracted data properly
    ingested_chunks = []
    for item in data:
        if 'text' in item and item['text'].strip():
            # Use the generator's ingest_documents function for consistency
            from app.services.generator import ingest_documents
            source_id = f"uploaded_{file.filename}_{str(uuid.uuid4())[:8]}"
            
            # For video/audio files, include additional metadata
            if 'word_count' in item and 'duration' in item:
                # Add metadata to the text for better context
                enhanced_text = f"[VIDEO/AUDIO TRANSCRIPT - Duration: {item['duration']:.2f}s, Words: {item['word_count']}]\n\n{item['text']}"
                chunk_ids, _ = ingest_documents([enhanced_text], source_id=source_id)
            else:
                chunk_ids, _ = ingest_documents([item['text']], source_id=source_id)
            
            ingested_chunks.extend(chunk_ids)
    
    return {
        "status": "parsed", 
        "filename": file.filename, 
        "data": data,
        "ingested_chunks": len(ingested_chunks),
        "message": f"Successfully ingested {len(ingested_chunks)} chunks from {file.filename}"
    }
