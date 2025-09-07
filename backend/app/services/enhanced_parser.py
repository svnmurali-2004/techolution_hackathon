"""
Enhanced Multi-File Parser with comprehensive file type support
"""
import os
import tempfile
import uuid
import requests
import pytesseract
from PIL import Image
import PyPDF2
import pandas as pd
from pptx import Presentation
from docx import Document
import zipfile
import tarfile
import rarfile
from chromadb import PersistentClient

# ChromaDB setup
CHROMA_PATH = "./chromadb_reports"
vector_db = PersistentClient(path=CHROMA_PATH)

# Deepgram API key
DEEPGRAM_API_KEY = "YOUR_DEEPGRAM_API_KEY"  # Replace with actual key

def extract_pdf(file_path):
    """Extract text from PDF files"""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += f"[Page {page_num + 1}]\n" + page.extract_text() + "\n\n"
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract PDF: {str(e)}"}]

def extract_pptx(file_path):
    """Extract text from PowerPoint presentations"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"[Slide {slide_num}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract PPTX: {str(e)}"}]

def extract_xlsx(file_path):
    """Extract text from Excel files"""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        text = ""
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text += f"[Sheet: {sheet_name}]\n"
            text += df.to_string(index=False) + "\n\n"
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract XLSX: {str(e)}"}]

def extract_txt(file_path):
    """Extract text from text files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract TXT: {str(e)}"}]

def extract_image(file_path):
    """Extract text from images using OCR"""
    try:
        text = pytesseract.image_to_string(file_path)
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract image text: {str(e)}"}]

def extract_docx(file_path):
    """Extract text from Word documents"""
    try:
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return [{"text": text.strip()}]
    except Exception as e:
        return [{"error": f"Failed to extract DOCX: {str(e)}"}]

def extract_video_audio(file_path):
    """Extract transcript from video/audio files using Deepgram API"""
    try:
        # Determine content type based on file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext in ['.wav', '.mp3', '.m4a', '.flac', '.aac', '.ogg']:
            content_type = "audio/wav" if file_ext == '.wav' else "audio/mp4"
        elif file_ext in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv']:
            content_type = "video/mp4"
        else:
            content_type = "video/mp4"  # Default for video files
        
        # Open the file and send to Deepgram
        with open(file_path, "rb") as f:
            response = requests.post(
                "https://api.deepgram.com/v1/listen?punctuate=true&diarize=true&model=nova-2",
                headers={
                    "Authorization": f"Token {DEEPGRAM_API_KEY}",
                    "Content-Type": content_type,
                },
                data=f,
                timeout=300  # 5 minute timeout for large files
            )
        
        if response.status_code != 200:
            raise Exception(f"Deepgram API error: {response.status_code} - {response.text}")
        
        result = response.json()
        
        # Extract transcript
        transcript = result["results"]["channels"][0]["alternatives"][0]["transcript"]
        
        # Extract word-level timestamps for better context
        words = result["results"]["channels"][0]["alternatives"][0]["words"]
        
        # Create structured data with timestamps
        import re
        sentences = re.split(r'[.!?]+', transcript)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        # Create chunks of 2-3 sentences each to avoid too many citations
        chunk_size = 2
        chunks = []
        for i in range(0, len(sentences), chunk_size):
            chunk_sentences = sentences[i:i + chunk_size]
            chunk_text = '. '.join(chunk_sentences)
            if chunk_text and not chunk_text.endswith('.'):
                chunk_text += '.'
            
            chunks.append({
                "text": chunk_text,
                "word_count": len(chunk_text.split()),
                "duration": words[-1]["end"] if words else 0,
                "words": words,
                "chunk_index": len(chunks) + 1
            })
        
        return chunks if chunks else [{
            "text": transcript,
            "word_count": len(words),
            "duration": words[-1]["end"] if words else 0,
            "words": words
        }]
        
    except Exception as e:
        print(f"Error transcribing video/audio: {str(e)}")
        return [{"error": f"Transcription failed: {str(e)}"}]

def extract_archive(file_path):
    """Extract text from archive files"""
    try:
        text_content = []
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == ".zip":
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                for file_info in zip_ref.filelist:
                    if not file_info.is_dir() and file_info.filename.endswith(('.txt', '.pdf', '.docx')):
                        try:
                            content = zip_ref.read(file_info.filename)
                            if file_info.filename.endswith('.txt'):
                                text_content.append(f"[From {file_info.filename}]: {content.decode('utf-8', errors='ignore')}")
                        except:
                            continue
        elif file_ext in [".tar", ".gz"]:
            with tarfile.open(file_path, 'r:*') as tar_ref:
                for member in tar_ref.getmembers():
                    if member.isfile() and member.name.endswith(('.txt', '.pdf', '.docx')):
                        try:
                            content = tar_ref.extractfile(member).read()
                            if member.name.endswith('.txt'):
                                text_content.append(f"[From {member.name}]: {content.decode('utf-8', errors='ignore')}")
                        except:
                            continue
        
        if text_content:
            return [{"text": "\n\n".join(text_content)}]
        else:
            return [{"text": f"Archive {os.path.basename(file_path)} contains no extractable text files"}]
            
    except Exception as e:
        return [{"error": f"Failed to extract archive: {str(e)}"}]

async def parse_file(file) -> dict:
    """Parse a single file and return extracted data"""
    ext = os.path.splitext(file.filename)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name
    
    # Handle different file types with comprehensive support
    if ext == ".pdf":
        data = extract_pdf(tmp_path)
    elif ext in [".ppt", ".pptx"]:
        data = extract_pptx(tmp_path)
    elif ext in [".xls", ".xlsx", ".csv"]:
        data = extract_xlsx(tmp_path)
    elif ext in [".txt", ".rtf"]:
        data = extract_txt(tmp_path)
    elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]:
        data = extract_image(tmp_path)
    elif ext in [".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv"]:
        data = extract_video_audio(tmp_path)
    elif ext in [".wav", ".mp3", ".m4a", ".flac", ".aac", ".ogg"]:
        data = extract_video_audio(tmp_path)
    elif ext in [".doc", ".docx"]:
        data = extract_docx(tmp_path)
    elif ext in [".zip", ".rar", ".7z", ".tar", ".gz"]:
        data = extract_archive(tmp_path)
    else:
        data = [{"error": f"Unsupported file type: {ext}"}]
    
    os.remove(tmp_path)
    
    # Store in vector DB using the same collection as generator
    collection = vector_db.get_or_create_collection("reports")
    
    # Process and store the extracted data properly
    ingested_chunks = []
    for item in data:
        if 'text' in item and item['text'].strip():
            # Use the generator's ingest_documents function for consistency
            from app.services.generator import ingest_documents
            source_id = f"uploaded_{file.filename}_{str(uuid.uuid4())[:8]}"
            
            # For video/audio files, include additional metadata
            if 'word_count' in item and 'duration' in item:
                # Add metadata to the text for better context
                enhanced_text = f"[VIDEO/AUDIO TRANSCRIPT - Duration: {item['duration']:.2f}s, Words: {item['word_count']}]\n\n{item['text']}"
                chunk_ids, _ = ingest_documents([enhanced_text], source_id=source_id)
            else:
                chunk_ids, _ = ingest_documents([item['text']], source_id=source_id)
            
            ingested_chunks.extend(chunk_ids)
    
    return {
        "status": "parsed", 
        "filename": file.filename,
        "chunks_ingested": len(ingested_chunks),
        "message": f"Successfully processed {file.filename}"
    }

async def parse_multiple_files(files) -> dict:
    """Parse multiple files and return comprehensive results"""
    results = {
        "total_files": len(files),
        "successful": 0,
        "failed": 0,
        "file_results": [],
        "total_chunks": 0
    }
    
    for file in files:
        try:
            result = await parse_file(file)
            if result["status"] == "parsed":
                results["successful"] += 1
                results["total_chunks"] += result.get("chunks_ingested", 0)
            else:
                results["failed"] += 1
            
            results["file_results"].append({
                "filename": file.filename,
                "status": result["status"],
                "message": result.get("message", "Unknown status"),
                "chunks": result.get("chunks_ingested", 0)
            })
            
        except Exception as e:
            results["failed"] += 1
            results["file_results"].append({
                "filename": file.filename,
                "status": "error",
                "message": str(e),
                "chunks": 0
            })
    
    return results

# File type categories for better organization
FILE_TYPE_CATEGORIES = {
    'documents': {
        'extensions': ['pdf', 'doc', 'docx', 'txt', 'rtf'],
        'description': 'Text documents and PDFs',
        'icon': '📄'
    },
    'presentations': {
        'extensions': ['ppt', 'pptx'],
        'description': 'PowerPoint presentations',
        'icon': '📊'
    },
    'spreadsheets': {
        'extensions': ['xls', 'xlsx', 'csv'],
        'description': 'Excel spreadsheets and CSV files',
        'icon': '📈'
    },
    'images': {
        'extensions': ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff'],
        'description': 'Images with OCR text extraction',
        'icon': '🖼️'
    },
    'videos': {
        'extensions': ['mp4', 'avi', 'mov', 'mkv', 'wmv', 'flv'],
        'description': 'Video files with transcript extraction',
        'icon': '🎥'
    },
    'audio': {
        'extensions': ['mp3', 'wav', 'm4a', 'flac', 'aac', 'ogg'],
        'description': 'Audio files with transcript extraction',
        'icon': '🎵'
    },
    'archives': {
        'extensions': ['zip', 'rar', '7z', 'tar', 'gz'],
        'description': 'Archive files with text extraction',
        'icon': '📦'
    }
}

def get_supported_file_types():
    """Get comprehensive list of supported file types"""
    all_extensions = []
    for category, info in FILE_TYPE_CATEGORIES.items():
        all_extensions.extend(info['extensions'])
    
    return {
        "categories": FILE_TYPE_CATEGORIES,
        "all_extensions": all_extensions,
        "total_types": len(all_extensions)
    }

def validate_file_type(filename):
    """Validate if file type is supported"""
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    supported_types = get_supported_file_types()["all_extensions"]
    return ext in supported_types, ext

def get_file_category(filename):
    """Get the category of a file based on its extension"""
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    for category, info in FILE_TYPE_CATEGORIES.items():
        if ext in info['extensions']:
            return category, info
    return "unknown", {"description": "Unknown file type", "icon": "❓"}
